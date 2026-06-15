"""
Gerador de Artigo (DOCX) e Apresentação (PPTX) — IME ET-261400
Idioma: Português Brasil
Execução: python gerar_artigo.py
"""
from docx import Document
from docx.shared import Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN
import os, datetime

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─── Paleta ─────────────────────────────────────────────────────────────────
C_NAVY  = RGBColor(0x1F, 0x38, 0x64)
C_BLUE  = RGBColor(0x2E, 0x75, 0xB6)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_GRAY  = RGBColor(0x70, 0x70, 0x70)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

P_NAVY  = PRGB(0x1F, 0x38, 0x64)
P_BLUE  = PRGB(0x2E, 0x75, 0xB6)
P_LBLUE = PRGB(0xBD, 0xD7, 0xEE)
P_WHITE = PRGB(0xFF, 0xFF, 0xFF)
P_BLACK = PRGB(0x00, 0x00, 0x00)
P_GRAY  = PRGB(0x40, 0x40, 0x40)
P_LGRAY = PRGB(0x9A, 0x9A, 0x9A)

RODAPE_PPTX = (
    "Gestão à Vista de OAEs Ferroviárias  |  IME — ET-261400  |  Wagner Tavares  |  2026"
)

# ════════════════════════════════════════════════════════════════════════════
#  Utilitários DOCX
# ════════════════════════════════════════════════════════════════════════════

def _font(run, name="Arial", size=12, bold=False, italic=False, color=None):
    run.font.name  = name
    run.font.size  = Pt(size)
    run.bold       = bold
    run.italic     = italic
    if color:
        run.font.color.rgb = color

def _fmt(p, before=0, after=6, ls=1.15, align=WD_ALIGN_PARAGRAPH.JUSTIFY):
    f = p.paragraph_format
    f.space_before      = Pt(before)
    f.space_after       = Pt(after)
    f.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    f.line_spacing      = ls
    p.alignment         = align

def _body(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    _font(r)
    _fmt(p)
    return p

def _h1(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    _font(r, size=13, bold=True, color=C_NAVY)
    _fmt(p, before=14, after=4, align=WD_ALIGN_PARAGRAPH.LEFT)
    return p

def _h2(doc, text):
    p = doc.add_paragraph()
    r = p.add_run(text)
    _font(r, size=12, bold=True, color=C_BLUE)
    _fmt(p, before=8, after=3, align=WD_ALIGN_PARAGRAPH.LEFT)
    return p

def _equation(doc, eq_text):
    p = doc.add_paragraph()
    r = p.add_run(eq_text)
    _font(r, size=12, italic=True)
    _fmt(p, before=4, after=4, align=WD_ALIGN_PARAGRAPH.CENTER)
    return p

def _set_cell_bg(cell, hex_color):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  hex_color)
    tcPr.append(shd)

def _cell_text(cell, text, bold=False, size=11, color=None,
               align=WD_ALIGN_PARAGRAPH.CENTER):
    p = cell.paragraphs[0]
    p.clear()
    r = p.add_run(text)
    _font(r, size=size, bold=bold, color=color or C_BLACK)
    _fmt(p, before=2, after=2, align=align)

def _fig_placeholder(doc, fig_num, caption_text, instrucao_captura):
    """Caixa cinza com instruções de captura + legenda da figura."""
    # Caixa de instrução de captura
    instr = doc.add_paragraph()
    instr_r = instr.add_run(f"✦ INSTRUÇÃO DE CAPTURA — FIGURA {fig_num}")
    _font(instr_r, size=10, bold=True, color=C_BLUE)
    _fmt(instr, before=10, after=2, align=WD_ALIGN_PARAGRAPH.LEFT)

    instr2 = doc.add_paragraph()
    instr2_r = instr2.add_run(instrucao_captura)
    _font(instr2_r, size=10, italic=True, color=C_GRAY)
    _fmt(instr2, before=0, after=4, align=WD_ALIGN_PARAGRAPH.LEFT)

    # Caixa placeholder cinza
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0)
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd  = OxmlElement('w:shd')
    shd.set(qn('w:val'),   'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'),  'D9D9D9')
    tcPr.append(shd)
    trPr = tc.getparent().get_or_add_trPr()
    trH  = OxmlElement('w:trHeight')
    trH.set(qn('w:val'), '2200')
    trPr.append(trH)
    cp  = cell.paragraphs[0]
    cpr = cp.add_run(f"[ INSERIR FIGURA {fig_num} AQUI — captura de tela do simulador ]")
    _font(cpr, size=11, italic=True, color=C_GRAY)
    _fmt(cp, before=8, after=8, align=WD_ALIGN_PARAGRAPH.CENTER)

    # Legenda
    doc.add_paragraph()
    cap   = doc.add_paragraph()
    cap_r = cap.add_run(f"Figura {fig_num}. {caption_text} Fonte: O autor (2026).")
    _font(cap_r, size=11, italic=True)
    _fmt(cap, before=2, after=12, align=WD_ALIGN_PARAGRAPH.CENTER)

def _ref_entry(doc, text):
    p   = doc.add_paragraph()
    run = p.add_run(text)
    _font(run, size=11)
    f = p.paragraph_format
    f.space_before      = Pt(0)
    f.space_after       = Pt(6)
    f.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    f.line_spacing      = 1.15
    f.left_indent       = Cm(1.25)
    f.first_line_indent = Cm(-1.25)
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    return p

# ════════════════════════════════════════════════════════════════════════════
#  BUILD DOCX
# ════════════════════════════════════════════════════════════════════════════

def build_docx():
    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Cm(2.5)
        sec.bottom_margin = Cm(2.5)
        sec.left_margin   = Cm(2.5)
        sec.right_margin  = Cm(2.5)

    # ── PÁGINA DE TÍTULO ────────────────────────────────────────────────────
    titulo = doc.add_paragraph()
    tr     = titulo.add_run(
        "Gestão à Vista de Obras de Arte Especiais Ferroviárias:\n"
        "Uma Abordagem Espacial e Orientada a Dados (Data-Centric)"
    )
    _font(tr, size=16, bold=True, color=C_NAVY)
    _fmt(titulo, before=24, after=12, align=WD_ALIGN_PARAGRAPH.CENTER)

    for linha in [
        "Wagner de Almeida Tavares¹",
        "¹ Instituto Militar de Engenharia (IME) — Programa de Pós-Graduação em Engenharia de Transportes\n"
        "  Rio de Janeiro, Brasil  |  wagner.tavares2024@gmail.com",
        f"Entregue em: {datetime.date.today().strftime('%d de %B de %Y')}",
    ]:
        p = doc.add_paragraph()
        r = p.add_run(linha)
        _font(r, size=11, italic=("wagner" in linha.lower() or "Entregue" in linha))
        _fmt(p, before=4, after=4, align=WD_ALIGN_PARAGRAPH.CENTER)

    doc.add_paragraph()

    # ── RESUMO ──────────────────────────────────────────────────────────────
    _h1(doc, "Resumo")
    _body(doc, (
        "A gestão da infraestrutura ferroviária depende de dados sistemáticos de inspeção de "
        "centenas de Obras de Arte Especiais (OAEs). No Brasil, esses dados permanecem fragmentados "
        "em documentos PDF, planilhas e sistemas desconectados, impedindo a identificação visual "
        "de vulnerabilidades ao longo dos corredores ferroviários. Este artigo apresenta um simulador "
        "de gestão à vista de OAEs ferroviárias, desenvolvido e validado com cinco anos de dados de "
        "inspeção estrutural (2021–2025) de um corredor de 905 km no eixo Minas Gerais–Espírito Santo. "
        "O sistema implementa: (i) pipeline ETL com lógica de fallback temporal para dados de inspeção "
        "multitemporais; (ii) algoritmo original de delimitação de zonas de influência do Risco de "
        "Agressividade Ambiental Estrutural (RAAE), combinando agrupamento Union-Find, Árvore Geradora "
        "Mínima de Prim e Iteração de Jacobi para produzir zonas mutuamente tangentes; "
        "(iii) cartografia web interativa multicamada com Folium/Leaflet.js; e (iv) geração "
        "automatizada de relatórios técnicos HTML alinhados à ABNT NBR 9452:2019. Validado com "
        "25 OAEs, o simulador elimina a fragmentação informacional, viabiliza a visualização de "
        "vulnerabilidades por trecho ferroviário e estabelece base geoespacial unificada para "
        "futuras aplicações de aprendizado profundo na previsão de degradação estrutural."
    ))
    kw   = doc.add_paragraph()
    kw_r = kw.add_run(
        "Palavras-chave: Obras de Arte Especiais; Gestão à Vista; RAAE; Inspeção Ferroviária; "
        "Árvore Geradora Mínima; Dashboard Interativo; Ciência de Dados."
    )
    _font(kw_r, size=11, italic=True)
    _fmt(kw, before=6, after=12, align=WD_ALIGN_PARAGRAPH.LEFT)

    # ── 1. INTRODUÇÃO ────────────────────────────────────────────────────────
    _h1(doc, "1.  Introdução")
    _body(doc, (
        "A malha ferroviária brasileira de cargas supera 30.000 km e inclui milhares de Obras de "
        "Arte Especiais (OAEs) — pontes, viadutos, passagens inferiores, túneis e passarelas — cuja "
        "integridade estrutural é condição essencial para a operação ininterrupta do transporte de "
        "cargas (ANTF, 2023). A norma ABNT NBR 9452:2019 estabelece três modalidades de inspeção "
        "obrigatórias: rotineira (anual), especial (a cada 60 meses) e instrumentada (sob condição "
        "específica). Cada ciclo gera notas de condição na escala 0 a 5 — onde 0 corresponde à "
        "condição Emergencial e 5 à condição Excelente — nas dimensões estrutural, de durabilidade "
        "e funcional, além de registros textuais de manifestações patológicas observadas."
    ))
    _body(doc, (
        "Apesar desse arcabouço normativo, os dados de inspeção de OAEs ferroviárias no Brasil "
        "permanecem sistematicamente fragmentados em documentos PDF, formulários Word e planilhas "
        "Excel desconectadas. Esse cenário impede a identificação visual de padrões de "
        "vulnerabilidade ao longo do corredor ferroviário e compromete o planejamento de manutenção "
        "baseado em evidências. Sistemas comerciais de gerenciamento de ativos existem como "
        "alternativas, mas impõem elevado custo de licenciamento e baixa adaptabilidade aos "
        "critérios das normas brasileiras de inspeção."
    ))
    _body(doc, (
        "O paradigma de inteligência artificial centrada em dados (Ng, 2021) reencuadra o desafio: "
        "em vez de construir algoritmos mais complexos, a prioridade é estruturar, limpar e "
        "governar os dados subjacentes — habilitando que modelos analíticos aprendam de entradas "
        "de alta qualidade. Essa abordagem é especialmente relevante para a gestão de OAEs "
        "ferroviárias, cujos conjuntos de dados são escassos, heterogêneos e historicamente "
        "subutilizados."
    ))
    _body(doc, (
        "A literatura recente confirma o potencial de métodos computacionais para transformar "
        "registros não estruturados de inspeção. Kang et al. (2020) demonstraram que redes neurais "
        "profundas extraem informações estruturadas de textos de inspeção de engenharia com alta "
        "fidelidade. Zheng et al. (2022) estenderam a abordagem para extração de relações semânticas "
        "em relatórios de inspeção de pontes, mapeando descrições textuais em grafos de dependência "
        "de falha estrutural. Para cenários com dados anotados escassos — comuns em contextos de "
        "inspeção de OAEs — Li et al. (2023) propuseram aprendizado few-shot aplicável a conjuntos "
        "piloto da escala aqui abordada. No plano sistêmico, Pozzi et al. (2017) demonstraram que "
        "a integração BIM-GIS potencializa substancialmente a capacidade analítica de sistemas de "
        "gerenciamento de pontes, motivando a arquitetura geoespacial adotada neste trabalho."
    ))
    _body(doc, (
        "Este artigo endereça a seguinte pergunta de pesquisa: ferramentas de ciência de dados "
        "de código aberto e algoritmos baseados em teoria dos grafos conseguem transformar dados "
        "fragmentados de inspeção de OAEs em um sistema de gestão à vista unificado e estruturado "
        "espacialmente? Os objetivos são: (1) projetar pipeline ETL para consolidação de dados de "
        "inspeção multitemporais; (2) desenvolver algoritmo original de delimitação automatizada "
        "de zonas RAAE; (3) entregar dashboard interativo multicamada para visualização de condição "
        "estrutural; e (4) validar o sistema completo em um corredor ferroviário real de 905 km."
    ))

    # ── 2. DADOS E MÉTODOS ───────────────────────────────────────────────────
    _h1(doc, "2.  Dados e Métodos")

    _h2(doc, "2.1  Área de Estudo e Fontes de Dados")
    _body(doc, (
        "O conjunto de dados de validação cobre um corredor ferroviário de 905 km no eixo "
        "Minas Gerais–Espírito Santo, composto por 25 OAEs na Fase 1. Três fontes primárias "
        "foram utilizadas: (i) dados espaciais — coordenadas GPS (datum WGS84) de cada OAE e "
        "arquivos KMZ de topologia ferroviária; (ii) registros de inspeção estrutural de "
        "2021 a 2025, em planilhas Excel multi-aba no formato ABNT NBR 9452:2019, abrangendo "
        "ciclos rotineiros e especiais com notas E, D, F e G; e (iii) descrições textuais de "
        "manifestações patológicas provenientes dos formulários de inspeção, cobrindo "
        "patologias em estruturas de concreto e aço."
    ))

    _h2(doc, "2.2  Pipeline ETL")
    _body(doc, (
        "O pipeline de Extração, Transformação e Carga (ETL) foi implementado em Python com "
        "Pandas (McKinney, 2010) e openpyxl. Três abas Excel — Dados Fixos, Inspeção Rotineira "
        "e Inspeção Especial — foram lidas por dicionários de índices posicionais de colunas, "
        "garantindo robustez a renomeações de cabeçalho. Um mecanismo de fallback temporal "
        "consulta os anos de 2025 a 2021 sequencialmente, selecionando o registro mais recente "
        "com nota geral válida. OAEs sem nenhum registro válido recebem o valor sentinela S.I. "
        "(sem inspeção), tratado como categoria distinta de dado ausente. As três tabelas são "
        "unidas por identificador de OAE após validação de coordenadas e normalização de strings."
    ))

    _h2(doc, "2.3  PLN Léxico para Extração de Patologias")
    _body(doc, (
        "A categorização automática de manifestações patológicas emprega PLN léxico em duas "
        "camadas: (i) dicionário de palavras-chave mapeando 10 categorias de patologias a "
        "expressões regulares — seis para concreto (armadura exposta, fissuração, carbonatação, "
        "segregação, deformação, calcinação) e quatro para aço (corrosão, fratura, deformação, "
        "falha em ligação) — aplicadas sem distinção de maiúsculas/minúsculas aos campos textuais; "
        "e (ii) coloração espacial: cada categoria ativa recebe uma cor distinta de uma paleta "
        "de 17 cores para visualização no mapa, permitindo leitura espacial da distribuição de "
        "patologias ao longo do corredor."
    ))

    _h2(doc, "2.4  Algoritmo de Delimitação de Zonas RAAE")
    _body(doc, (
        "A delimitação das zonas de influência do Risco de Agressividade Ambiental Estrutural "
        "(RAAE) — contribuição algorítmica central deste trabalho — resolve um problema de "
        "partição espacial unidimensional: atribuir a cada OAE um raio de influência r_i tal "
        "que círculos de OAEs vizinhas sejam mutuamente tangentes (r_i + r_j = d_ij), sem "
        "sobreposições e sem raios nulos. A solução ocorre em três etapas encadeadas."
    ))
    _body(doc, (
        "Etapa 1 — Agrupamento por ε-proximidade (Union-Find): OAEs dentro de distância geodésica "
        "ε = 500 m são agrupadas em um nó representativo no centroide geométrico do grupo, "
        "usando Union-Find com compressão de caminho (O(n·α(n)) amortizado). Isso trata estruturas "
        "co-localizadas que compartilham a mesma classe de agressividade ambiental."
    ))
    _body(doc, (
        "Etapa 2 — Árvore Geradora Mínima (Algoritmo de Prim): O algoritmo de Prim constrói a "
        "AGM sobre os k grupos resultantes usando fila de prioridade (O(k² log k)), definindo "
        "a topologia de conectividade que captura a estrutura linear da ferrovia sem requerer a "
        "geometria explícita do traçado como entrada. Essa adjacência corresponde funcionalmente "
        "a uma decomposição de Voronoi 1D ao longo do corredor (Aurenhammer, 1991)."
    ))
    _body(doc, (
        "Etapa 3 — Iteração de Jacobi: O sistema de restrições de tangência é resolvido "
        "iterativamente por relaxação simultânea. Para cada nó i com vizinhos N(i) na AGM:"
    ))
    _equation(doc, "r_i^(t+1) = max( r_min,  min_{j ∈ N(i)} ( d_ij − r_j^(t) ) )          (1)")
    _body(doc, (
        "onde r_min = 800 m e a convergência é declarada quando a variação máxima é inferior "
        "a 1 m, com fator de buffer de 3% para tangência visualmente perceptível na escala "
        "cartográfica. Arestas AGM com d_ij − r_i − r_j > 1.000 m são representadas como "
        "linhas tracejadas no mapa, identificando explicitamente trechos ferroviários sem "
        "cobertura de monitoramento."
    ))

    _h2(doc, "2.5  Arquitetura do Dashboard Interativo")
    _body(doc, (
        "O mapa interativo é construído com Folium (wrapper Python para Leaflet.js) com três "
        "camadas de tiles intercambiáveis: CartoDB Positron, OpenStreetMap e CartoDB Dark Matter. "
        "Distâncias são calculadas geodesicamente com GeoPy (elipsoide WGS84), garantindo "
        "precisão métrica independente da latitude. O dashboard é implementado em Streamlit "
        "com arquitetura modular em quatro camadas (data_loader → map_builder → app → "
        "report_generator), com filtros hierárquicos por tipo de OAE, característica estrutural, "
        "zona RAAE, tipo de patologia e ano de construção (Figura 1)."
    ))
    _fig_placeholder(doc, 1,
        "Visão geral do dashboard interativo: painel de KPIs de condição estrutural (escala "
        "ABNT NBR 9452:2019), barra lateral com filtros hierárquicos, mapa geoespacial "
        "multicamada e legenda RAAE com contagem por zona. ",
        "COMO CAPTURAR:\n"
        "1. Abrir o simulador em http://localhost:8501\n"
        "2. Barra lateral: manter todos os filtros no padrão (sem filtro ativo)\n"
        "3. Mapa: selecionar tile CartoDB Positron | Coloração = Por Condição\n"
        "4. Garantir que o painel de KPIs (Total OAEs, Com Rotineira, Com Especial) esteja visível\n"
        "5. Expandir a janela do navegador para resolução mínima 1280×900 px\n"
        "6. Capturar a tela inteira do navegador (Alt+Print Screen ou ferramenta Snipping Tool)\n"
        "7. Resolução mínima: 1200×800 px — formato JPEG ou PNG"
    )

    _h2(doc, "2.6  Geração Automatizada de Relatório Técnico")
    _body(doc, (
        "O módulo de relatório implementa o padrão Dado-para-Documento: a partir dos registros "
        "de inspeção filtrados pelo usuário, gera automaticamente um documento HTML técnico "
        "estruturado — exportável como PDF pela função de impressão do navegador — contendo "
        "capa com metadados da emissão, tabelas de resumo executivo e fichas individuais de "
        "OAE com notas de condição, manifestações patológicas e recomendações de intervenção "
        "vinculadas automaticamente à ABNT NBR 9452:2019, NBR 6118:2014 e NBR 12655."
    ))

    # ── 3. RESULTADOS ────────────────────────────────────────────────────────
    _h1(doc, "3.  Resultados")

    _h2(doc, "3.1  Consolidação ETL")
    _body(doc, (
        "O pipeline consolidou com sucesso 25 OAEs ao longo de cinco anos de inspeção. O "
        "fallback temporal resolveu 14 OAEs com registros de 2025, 7 com registros de 2024 e "
        "4 com registros de 2023, sem nenhuma OAE requerendo fallback além de 2023. A "
        "distribuição final de condição revelou: 0 OAEs em Emergencial (0) ou Crítico (1), "
        "4 em Ruim (2), 13 em Regular (3), 1 em Bom (4), 0 em Excelente (5) e 7 classificadas "
        "como S.I. (ciclo de inspeção especial pendente)."
    ))

    _h2(doc, "3.2  Distribuição RAAE e Desempenho do Algoritmo")
    _body(doc, (
        "A Tabela 1 resume a distribuição das zonas RAAE entre as 25 OAEs. A etapa de "
        "agrupamento Union-Find identificou 6 pares co-localizados, reduzindo as 25 OAEs "
        "a 19 grupos representativos. A AGM de Prim foi construída com 18 arestas sobre "
        "19 nós. A iteração de Jacobi convergiu em menos de 86 iterações para todos os grupos. "
        "Das 18 arestas AGM, 16 (88,9%) produziram zonas RAAE mutuamente tangentes. Duas "
        "arestas com gaps residuais de 35,9 km e 171,2 km — correspondentes a segmentos "
        "ferroviários sem OAEs intermediárias — são representadas como linhas tracejadas "
        "cinzas (Figura 2), convertendo a lacuna de dados em diagnóstico de cobertura de "
        "monitoramento."
    ))

    # Tabela 1
    doc.add_paragraph()
    tbl = doc.add_table(rows=6, cols=4)
    tbl.style     = 'Table Grid'
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    cabecalhos = ["Nível RAAE", "Qtd. OAEs", "Percentual", "Cor no Mapa"]
    for j, h in enumerate(cabecalhos):
        _set_cell_bg(tbl.cell(0, j), '1F3864')
        _cell_text(tbl.cell(0, j), h, bold=True, color=C_WHITE)
    linhas = [
        ("Baixo",      "7",  "28%", "Verde"),
        ("Médio",      "12", "48%", "Amarelo"),
        ("Alto",       "4",  "16%", "Laranja"),
        ("Muito Alto", "2",  "8%",  "Vermelho"),
        ("Extremo",    "0",  "0%",  "Roxo"),
    ]
    for i, (lvl, cnt, pct, cor) in enumerate(linhas):
        bg = 'D9E2F3' if i % 2 == 0 else 'FFFFFF'
        for j, val in enumerate([lvl, cnt, pct, cor]):
            _set_cell_bg(tbl.cell(i+1, j), bg)
            _cell_text(tbl.cell(i+1, j), val,
                       align=(WD_ALIGN_PARAGRAPH.LEFT if j == 0
                               else WD_ALIGN_PARAGRAPH.CENTER))
    cap_tbl = doc.add_paragraph()
    cap_r   = cap_tbl.add_run(
        "Tabela 1. Distribuição de zonas RAAE — 25 OAEs, corredor ferroviário MG-ES. "
        "Fonte: O autor (2026)."
    )
    _font(cap_r, size=11, italic=True)
    _fmt(cap_tbl, before=4, after=10, align=WD_ALIGN_PARAGRAPH.CENTER)

    _fig_placeholder(doc, 2,
        "Zonas de influência RAAE mutuamente tangentes ao longo do corredor ferroviário "
        "MG-ES (905 km, 25 OAEs), geradas pelo algoritmo Union-Find + AGM (Prim) + "
        "Iteração de Jacobi. Linhas tracejadas indicam trechos sem cobertura de "
        "monitoramento (35,9 km e 171,2 km). Camada base: CartoDB Positron. ",
        "COMO CAPTURAR:\n"
        "1. Abrir o simulador em http://localhost:8501\n"
        "2. Barra lateral: manter todos os filtros padrão\n"
        "3. Mapa: selecionar tile CartoDB Positron (fundo branco — melhor para impressão)\n"
        "4. Coloração = Por RAAE (botão de rádio na barra lateral)\n"
        "5. Ativar a camada 'Zonas RAAE' no controle de camadas do mapa (canto superior direito)\n"
        "6. Ajustar o zoom para mostrar TODO o corredor MG-ES com todas as 25 OAEs visíveis\n"
        "7. Verificar que os círculos coloridos são visíveis e a legenda RAAE aparece\n"
        "8. Capturar SOMENTE a área do mapa (sem a barra lateral)\n"
        "9. Resolução mínima: 1400×900 px — formato PNG (melhor qualidade)"
    )

    _h2(doc, "3.3  Dashboard Interativo")
    _body(doc, (
        "O dashboard Streamlit (Figura 3) integra o mapa Folium a painéis de KPIs em tempo "
        "real exibindo: total de OAEs (25), OAEs com inspeção rotineira (25), OAEs com "
        "inspeção especial (18), histograma de distribuição de condição e legenda de zonas "
        "RAAE com contagem por nível. O LayerControl permite ativação independente de até "
        "10 camadas de patologias, viabilizando análise de padrões de distribuição ao longo "
        "do corredor. A troca de tile (CartoDB Positron ↔ OpenStreetMap ↔ Dark Matter) "
        "não requer reconstrução do mapa, fornecendo visualização adaptada a cenários de "
        "inspeção, relatório e comunicação de risco."
    ))
    _fig_placeholder(doc, 3,
        "Dashboard interativo de gestão à vista com painéis de KPIs de condição (escala "
        "ABNT NBR 9452:2019), filtros hierárquicos na barra lateral e mapa multicamada "
        "com marcadores de OAE coloridos por zona RAAE e camadas de manifestações "
        "patológicas ativas. ",
        "COMO CAPTURAR:\n"
        "1. Abrir o simulador em http://localhost:8501\n"
        "2. Barra lateral: ativar pelo menos 2 ou 3 filtros para mostrar a funcionalidade\n"
        "   (ex.: filtrar por RAAE = Médio + Alto)\n"
        "3. Mapa: tile CartoDB Positron | Coloração = Por RAAE\n"
        "4. Ativar UMA camada de patologia (ex.: 'Armadura exposta') para mostrar o LayerControl\n"
        "5. Zoom no mapa: nível intermediário, mostrando 8–12 OAEs com círculos RAAE visíveis\n"
        "6. Garantir que a legenda RAAE (com contagens) esteja visível no canto inferior direito\n"
        "7. Capturar a tela COMPLETA do navegador (barra lateral + mapa + KPIs)\n"
        "8. Resolução mínima: 1280×900 px — formato PNG"
    )

    _h2(doc, "3.4  Geração Automatizada de Relatório")
    _body(doc, (
        "A seleção dos parâmetros de condição 2 (Ruim) e 3 (Regular) gerou relatório técnico "
        "para 17 OAEs em menos de dois segundos, incluindo: capa com metadados da emissão, "
        "tabelas de resumo executivo (distribuição de condição e distribuição RAAE), fichas "
        "individuais de OAE com notas de inspeção rotineira e especial, manifestações "
        "patológicas e recomendações de intervenção normativas. O processo anteriormente "
        "demandava elaboração manual por engenheiro responsável; o sistema automatiza e "
        "padroniza integralmente esse fluxo para qualquer seleção de condição."
    ))

    # ── 4. CONCLUSÕES ────────────────────────────────────────────────────────
    _h1(doc, "4.  Conclusões")
    _body(doc, (
        "Este trabalho desenvolveu e validou um simulador de gestão à vista de OAEs ferroviárias "
        "com ferramentas Python de código aberto, demonstrando que algoritmos baseados em teoria "
        "dos grafos integrados à visualização geoespacial resolvem a fragmentação informacional "
        "característica dos sistemas atuais de gestão de inspeções. Quatro contribuições "
        "principais foram entregues."
    ))
    _body(doc, (
        "Primeiro, um pipeline ETL multitemporal com lógica de fallback e tratamento de valor "
        "sentinela consolida registros heterogêneos em um DataFrame geoespacial único, pronto "
        "para análise. Segundo, o algoritmo original de delimitação de zonas RAAE em três etapas "
        "(Union-Find + AGM de Prim + Iteração de Jacobi) produz partição espacial equivalente "
        "a Voronoi 1D ao longo do corredor ferroviário sem requerer a geometria explícita do "
        "traçado. A identificação de dois segmentos não monitorados (35,9 km e 171,2 km) "
        "converte uma lacuna de dados em recomendação de cobertura de monitoramento. Terceiro, "
        "o dashboard interativo multicamada fornece visualização de condição estrutural para "
        "25 OAEs ao longo de 905 km. Quarto, a geração automatizada de relatório HTML reduz "
        "o esforço de elaboração manual de horas para segundos."
    ))
    _body(doc, (
        "Limitações incluem: o piloto cobre apenas 25 das 389 OAEs do corredor, limitando "
        "a generalização estatística; o PLN léxico não realiza desambiguação contextual para "
        "patologias compostas; e a iteração de Jacobi não garante otimalidade global, embora "
        "a convergência empírica seja estável para topologias de grafo quase-lineares."
    ))
    _body(doc, (
        "A arquitetura modular escala diretamente para a Fase 2 (389 OAEs) sem alteração de "
        "código. Direções futuras incluem: (1) expansão para o corredor completo; (2) integração "
        "de registros fotográficos com modelos de legendagem de imagens (Brilakis et al., 2023); "
        "e (3) treinamento de modelos LSTM ou Transformer para previsão probabilística de "
        "degradação estrutural, tendo como base o banco de dados de inspeção multianual "
        "estruturado por este pipeline — direção validada por Kang et al. (2020) e "
        "Zheng et al. (2022) em contextos análogos."
    ))

    # ── AGRADECIMENTOS ───────────────────────────────────────────────────────
    _h1(doc, "Agradecimentos")
    _body(doc, (
        "O autor agradece ao Instituto Militar de Engenharia (IME) e aos Professores Daniel "
        "Rodrigues dos Santos, Filipe Almeida Corrêa do Nascimento e Tiago Barreto Tamagusko "
        "pelo suporte acadêmico e orientação ao longo da disciplina ET-261400."
    ))

    # ── REFERÊNCIAS ──────────────────────────────────────────────────────────
    _h1(doc, "Referências")
    refs = [
        ("Associação Brasileira de Normas Técnicas. (2019). "
         "NBR 9452: Inspeção de pontes, viadutos e passagens subterrâneas [Norma]. ABNT."),
        ("Associação Brasileira de Normas Técnicas. (2014). "
         "NBR 6118: Projeto de estruturas de concreto — Procedimento [Norma]. ABNT."),
        ("Aurenhammer, F. (1991). Voronoi diagrams — A survey of a fundamental geometric "
         "data structure. ACM Computing Surveys, 23(3), 345–405. "
         "https://doi.org/10.1145/116873.116880"),
        ("Brilakis, I., Pan, Y., Borrmann, A., Doukari, O., & Siu, M. (2023). Image captioning "
         "for automated bridge inspection: A feasibility study. Automation in Construction, 148, "
         "104764. https://doi.org/10.1016/j.autcon.2023.104764"),
        ("Kang, J. S., Lee, H., Park, J., & Yoo, S. (2020). A deep neural network-based method "
         "for information extraction using transfer learning strategies to support automated "
         "compliance checking. Automation in Construction, 110, 103035. "
         "https://doi.org/10.1016/j.autcon.2020.103035"),
        ("Li, Y., Zhang, M., & Wang, Q. (2023). Few-shot learning for key information extraction "
         "from bridge inspection reports. Journal of Bridge Engineering, 28(4), 04023012. "
         "https://doi.org/10.1061/JBENF2.BEENG-5898"),
        ("McKinney, W. (2010). Data structures for statistical computing in Python. Em "
         "Proceedings of the 9th Python in Science Conference (pp. 51–56)."),
        ("Ng, A. (2021). A chat with Andrew on MLOps: From model-centric to data-centric AI "
         "[Vídeo]. DeepLearning.AI. https://www.deeplearning.ai"),
        ("Pozzi, M., Der Kiureghian, A., & Castelletti, A. (2017). Development of a bridge "
         "management system integrating BIM/GIS. Journal of Bridge Engineering, 22(11), 04017083. "
         "https://doi.org/10.1061/(ASCE)BE.1943-5592.0001107"),
        ("Streamlit Inc. (2019). Streamlit: The fastest way to build and share data apps "
         "[Software]. https://streamlit.io"),
        ("Zheng, X., Dong, J., & Wu, C. (2022). Semantic neural network ensemble for dependency "
         "relation extraction from bridge inspection reports. Computer-Aided Civil and "
         "Infrastructure Engineering, 37(8), 987–1005. "
         "https://doi.org/10.1111/mice.12793"),
    ]
    for ref in refs:
        _ref_entry(doc, ref)

    out = os.path.join(OUT_DIR, "artigo_oae_ET261400_ptBR.docx")
    doc.save(out)
    print(f"[DOCX] Salvo: {out}")
    return out


# ════════════════════════════════════════════════════════════════════════════
#  Utilitários PPTX
# ════════════════════════════════════════════════════════════════════════════

def _slide_bg(slide, hex_color):
    bg   = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = PRGB(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def _box(slide, left, top, w, h, text, fs=17, bold=False, italic=False,
         color=P_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(PIn(left), PIn(top), PIn(w), PIn(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.size      = PPt(fs)
    r.font.color.rgb = color
    return tb

def _rect(slide, left, top, w, h, fill_hex, line_hex=None):
    sh = slide.shapes.add_shape(
        1, PIn(left), PIn(top), PIn(w), PIn(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PRGB(
        int(fill_hex[0:2],16), int(fill_hex[2:4],16), int(fill_hex[4:6],16))
    if line_hex:
        sh.line.color.rgb = PRGB(
            int(line_hex[0:2],16), int(line_hex[2:4],16), int(line_hex[4:6],16))
    else:
        sh.line.fill.background()
    return sh

def _rodape(slide):
    _box(slide, 0, 6.88, 10, 0.22, RODAPE_PPTX,
         fs=8.5, color=P_LGRAY, align=PP_ALIGN.CENTER)

def _slide_titulo(prs, titulo, subtitulo):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, '1F3864')
    _rect(sl, 0, 3.1, 10, 0.07, '2E75B6')
    _box(sl, 0.7, 0.6, 8.6, 2.2, titulo,
         fs=28, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER, wrap=True)
    _box(sl, 0.7, 3.3, 8.6, 1.0, subtitulo,
         fs=17, italic=True, color=P_LBLUE, align=PP_ALIGN.CENTER, wrap=True)
    _box(sl, 0.7, 4.5, 8.6, 0.55,
         "Wagner de Almeida Tavares  |  IME — ET-261400  |  Junho 2026",
         fs=13, color=P_LBLUE, align=PP_ALIGN.CENTER)
    return sl

def _slide_divisor(prs, num, titulo_secao):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, '2E75B6')
    _box(sl, 0.5, 1.7, 9.0, 0.6, f"Seção {num}",
         fs=22, color=P_LBLUE, align=PP_ALIGN.CENTER)
    _box(sl, 0.5, 2.4, 9.0, 1.3, titulo_secao,
         fs=34, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER, wrap=True)
    return sl

def _slide_conteudo(prs, titulo, bullets):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, 'FFFFFF')
    _rect(sl, 0, 0, 10, 0.95, '1F3864')
    _box(sl, 0.25, 0.10, 9.5, 0.78, titulo,
         fs=21, bold=True, color=P_WHITE, align=PP_ALIGN.LEFT)
    if isinstance(bullets, str):
        bullets = [bullets]
    y = 1.10
    for b in bullets:
        sub = b.startswith("  ")
        fs  = 14.5 if not sub else 13
        indent = 0.30 if not sub else 0.60
        pref   = "▸  " if not sub else "    —  "
        col    = P_BLACK if not sub else P_GRAY
        _box(sl, indent, y, 9.35 - indent, 0.46,
             pref + b.strip(), fs=fs, color=col, wrap=True)
        y += 0.44 if not sub else 0.38
    _rodape(sl)
    return sl

def _slide_duas_colunas(prs, titulo, col_esq, col_dir):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, 'FFFFFF')
    _rect(sl, 0, 0, 10, 0.95, '1F3864')
    _box(sl, 0.25, 0.10, 9.5, 0.78, titulo,
         fs=21, bold=True, color=P_WHITE)
    _rect(sl, 0.20, 1.05, 4.50, 0.40, '2E75B6')
    _box(sl, 0.20, 1.05, 4.50, 0.40, col_esq[0],
         fs=13, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    y = 1.55
    for b in col_esq[1:]:
        _box(sl, 0.30, y, 4.25, 0.42, "▸  " + b, fs=13, color=P_BLACK, wrap=True)
        y += 0.40
    _rect(sl, 5.10, 1.05, 4.65, 0.40, '1F3864')
    _box(sl, 5.10, 1.05, 4.65, 0.40, col_dir[0],
         fs=13, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    y = 1.55
    for b in col_dir[1:]:
        _box(sl, 5.20, y, 4.40, 0.42, "▸  " + b, fs=13, color=P_BLACK, wrap=True)
        y += 0.40
    _rodape(sl)
    return sl

def _slide_figura(prs, titulo, legenda, instrucao):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, 'FFFFFF')
    _rect(sl, 0, 0, 10, 0.95, '1F3864')
    _box(sl, 0.25, 0.10, 9.5, 0.78, titulo,
         fs=21, bold=True, color=P_WHITE)
    # Caixa placeholder
    _rect(sl, 0.40, 1.05, 9.20, 4.30, 'E8E8E8', 'BFBFBF')
    _box(sl, 0.40, 1.05, 9.20, 3.60,
         "📷  " + instrucao,
         fs=13, italic=True, color=PRGB(0x60,0x60,0x60),
         align=PP_ALIGN.CENTER, wrap=True)
    _box(sl, 0.40, 5.42, 9.20, 0.48,
         legenda, fs=10.5, italic=True, color=P_GRAY,
         align=PP_ALIGN.CENTER, wrap=True)
    _rodape(sl)
    return sl

def _slide_tabela(prs, titulo, cabecalhos, linhas, nota=""):
    from pptx.util import Inches as I2, Pt as P2
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, 'FFFFFF')
    _rect(sl, 0, 0, 10, 0.95, '1F3864')
    _box(sl, 0.25, 0.10, 9.5, 0.78, titulo,
         fs=21, bold=True, color=P_WHITE)
    nc = len(cabecalhos)
    nr = len(linhas)
    tbl = sl.shapes.add_table(
        nr+1, nc, PIn(0.60), PIn(1.15), PIn(8.80), PIn(nr*0.52+0.52)).table
    for j, h in enumerate(cabecalhos):
        c = tbl.cell(0, j)
        c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = P_NAVY
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if p.runs:
            p.runs[0].font.bold  = True
            p.runs[0].font.size  = PPt(13)
            p.runs[0].font.color.rgb = P_WHITE
    for i, row in enumerate(linhas):
        bg = PRGB(0xBD,0xD7,0xEE) if i%2==0 else PRGB(0xFF,0xFF,0xFF)
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j)
            c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.size = PPt(13)
    if nota:
        _box(sl, 0.60, 1.15 + nr*0.52+0.52 + 0.06, 8.80, 0.36,
             nota, fs=10, italic=True, color=P_GRAY, align=PP_ALIGN.CENTER)
    _rodape(sl)
    return sl

# ════════════════════════════════════════════════════════════════════════════
#  BUILD PPTX
# ════════════════════════════════════════════════════════════════════════════

def build_pptx():
    prs = Presentation()
    prs.slide_width  = PIn(10)
    prs.slide_height = PIn(7.5)

    # 1 ── Capa
    _slide_titulo(prs,
        "Gestão à Vista de OAEs Ferroviárias:\nUma Abordagem Espacial e Orientada a Dados",
        "ET-261400 — Ciência de Dados e Aprendizado Profundo aplicados aos Transportes")

    # 2 ── Roteiro
    _slide_conteudo(prs, "Roteiro da Apresentação", [
        "1.  Introdução — Problema, Contexto e Pergunta de Pesquisa",
        "2.  Referencial — Ciência de Dados Centrada em Dados e Trabalhos Relacionados",
        "3.  Dados e Métodos — ETL, PLN, Algoritmo RAAE, Dashboard",
        "4.  Resultados — Mapas, Tabelas e Relatório Automatizado",
        "5.  Conclusões — Contribuições, Limitações e Trabalhos Futuros",
    ])

    # ─── SEÇÃO 1: INTRODUÇÃO ─────────────────────────────────────────────────
    _slide_divisor(prs, 1, "Introdução")

    # 4
    _slide_conteudo(prs, "Problema: Fragmentação dos Dados de Inspeção de OAEs", [
        "Malha ferroviária brasileira: > 30.000 km com centenas de OAEs",
        "OAEs = pontes, viadutos, passagens inferiores, túneis, passarelas",
        "ABNT NBR 9452:2019 exige inspeções estruturadas — mas os dados ficam isolados:",
        "  Relatórios PDF e formulários Word — não processáveis por máquina",
        "  Planilhas Excel desconectadas — formatos heterogêneos por ciclo",
        "  Sem referenciamento geoespacial — impossível ver padrões no corredor",
        "Resultado: Nenhuma visão sistêmica e visual de vulnerabilidade ferroviária",
        "Sistemas comerciais GIS/BMS: alto custo, baixa adaptabilidade às normas brasileiras",
    ])

    # 5
    _slide_conteudo(prs, "Contexto: Corredor Ferroviário MG-ES", [
        "Corredor estudado: 905 km — Minas Gerais → Espírito Santo",
        "Fase 1 (piloto): 25 OAEs instrumentadas — pontes, viadutos, passagens inferiores",
        "Fase 2 (meta): 389 OAEs no mesmo corredor",
        "Ciclos de inspeção (NBR 9452:2019):",
        "  Rotineira — anual | escala de condição 0 (Emergencial) a 5 (Excelente)",
        "  Especial — a cada 60 meses | avaliação estrutural detalhada",
        "  Instrumentada — sob condição específica",
        "Notas por dimensão: E (estrutural) · D (durabilidade) · F (funcional) · G (geral)",
    ])

    # 6
    _slide_conteudo(prs, "Escala de Condição ABNT NBR 9452:2019", [
        "0  Emergencial  — Interdição imediata; risco estrutural crítico",
        "1  Crítico       — Intervenção em até 30 dias",
        "2  Ruim          — Intervenção em até 90 dias",
        "3  Regular       — Intervenção em até 180 dias",
        "4  Bom           — Ciclo anual de manutenção",
        "5  Excelente     — Plano de manutenção preventiva",
        "RAAE — Risco de Agressividade Ambiental Estrutural:",
        "  Baixo · Médio · Alto · Muito Alto · Extremo",
    ])

    # 7
    _slide_conteudo(prs, "Paradigma Data-Centric (Ng, 2021)", [
        "Abordagem tradicional: melhorar o modelo → mais acurácia",
        "Abordagem centrada em dados: melhorar os dados → mais acurácia em qualquer modelo",
        "Por que é relevante para OAEs ferroviárias?",
        "  Dados escassos (25–389 estruturas) — insuficientes para abordagens model-centric",
        "  Formatos heterogêneos — exigem ETL estruturado antes de qualquer análise",
        "  Subutilização histórica — anos de dados nunca integrados sistematicamente",
        "Este trabalho: construir a infraestrutura de dados primeiro → habilitar DNNs futuras",
    ])

    # 8
    _slide_conteudo(prs, "Pergunta de Pesquisa e Objetivos", [
        "Pergunta de Pesquisa:",
        "  Ferramentas de código aberto e algoritmos de teoria dos grafos conseguem transformar",
        "  dados fragmentados de inspeção de OAEs em um sistema de gestão à vista unificado?",
        "Objetivos:",
        "  (1) Pipeline ETL para consolidação multitemporal de dados de inspeção",
        "  (2) Algoritmo original de delimitação de zonas RAAE (grafos + otimização)",
        "  (3) Dashboard interativo multicamada para visualização de condição estrutural",
        "  (4) Validação em corredor ferroviário real de 905 km",
    ])

    # ─── SEÇÃO 2: REFERENCIAL ────────────────────────────────────────────────
    _slide_divisor(prs, 2, "Referencial Teórico e Trabalhos Relacionados")

    # 10
    _slide_duas_colunas(prs, "PLN Aplicado à Inspeção de Pontes",
        ["Extração de Texto (Kang et al., 2020)",
         "Redes neurais profundas extraem informações estruturadas de laudos de engenharia",
         "Transfer learning para relatórios heterogêneos",
         "Valida nossa abordagem de dicionário de termos + regex",
         "Few-Shot Learning (Li et al., 2023)",
         "Extração com poucos dados anotados",
         "Aplicável ao piloto de 25 OAEs"],
        ["PLN Semântico (Zheng et al., 2022)",
         "Ensemble neural para extração de relações de dependência",
         "Mapeia texto → grafos de falha estrutural",
         "Direção futura para nossos dados de patologia",
         "Legendagem de Imagens (Brilakis et al., 2023)",
         "Relatórios automáticos a partir de fotos de inspeção",
         "Alvo de integração na Fase 2"])

    # 11
    _slide_conteudo(prs, "BMS com GIS e Lacuna de Pesquisa", [
        "Pozzi et al. (2017): integração BIM-GIS potencializa sistemas de gestão de pontes",
        "  Motivação direta para a arquitetura geoespacial adotada neste trabalho",
        "Lacuna identificada:",
        "  Sistemas BMS proprietários: caros, pouco adaptáveis às normas brasileiras",
        "  Sistemas NLP acadêmicos: processam texto, mas sem camada de visualização geoespacial",
        "  Não existe sistema de código aberto que integre:",
        "  ETL + algoritmos geoespaciais + dashboard interativo + contexto NBR 9452",
        "Este trabalho preenche essa lacuna usando exclusivamente ferramentas Python open-source",
    ])

    # ─── SEÇÃO 3: DADOS E MÉTODOS ────────────────────────────────────────────
    _slide_divisor(prs, 3, "Dados e Métodos")

    # 13
    _slide_conteudo(prs, "Área de Estudo e Fontes de Dados", [
        "Área de estudo: corredor ferroviário de 905 km — MG → ES, Brasil",
        "Fase 1: 25 OAEs (pontes, viadutos, passagens inferiores)",
        "Fonte 1 — Dados Espaciais:",
        "  Coordenadas GPS (WGS84) de cada OAE · Arquivos KMZ de topologia",
        "Fonte 2 — Registros de Inspeção Estrutural (2021–2025):",
        "  Planilhas Excel multi-aba (formato NBR 9452:2019)",
        "  Notas E, D, F, G — inspeções rotineiras e especiais",
        "Fonte 3 — Textos de Manifestações Patológicas:",
        "  Campos de texto livre dos formulários de inspeção",
        "  Patologias em concreto e aço",
    ])

    # 14
    _slide_conteudo(prs, "Pipeline ETL — Arquitetura e Lógica de Fallback", [
        "Ferramentas: Python · Pandas · openpyxl",
        "Extração: 3 abas Excel via dicionários de índices posicionais de colunas",
        "  Robusto a renomeações de cabeçalho — sem nomes de colunas fixos no código",
        "Lógica de Fallback Temporal:",
        "  Busca anos: 2025 → 2024 → 2023 → 2022 → 2021",
        "  Seleciona registro mais recente com nota geral (G) válida",
        "  OAEs sem nenhum registro → sentinela S.I. (sem inspeção)",
        "Transformação: normalização de strings, validação de coordenadas, parsing de ano",
        "Carga: DataFrame Pandas com colunas de geometria (lat, lon) — pronto para análise",
    ])

    # 15
    _slide_conteudo(prs, "PLN Léxico — Extração de Manifestações Patológicas", [
        "PLN léxico em 2 camadas sobre campos de texto livre de inspeção:",
        "Camada 1 — Dicionário de Termos (regex, sem distinção maiúsc./minúsc.):",
        "  Concreto: armadura exposta · fissuração · carbonatação · segregação · deformação · calcinação",
        "  Aço: corrosão · fratura · deformação · falha em ligação",
        "  → 10 categorias de patologia · paleta de 17 cores",
        "Camada 2 — Visualização Espacial:",
        "  Cada tipo de patologia ativa → camada independente no mapa (LayerControl Folium)",
        "  Leitura espacial da distribuição de patologias ao longo do corredor",
        "Evolução futura: PLN semântico (Zheng et al., 2022) para Fase 2",
    ])

    # 16
    _slide_conteudo(prs, "Arquitetura do Sistema — 4 Camadas", [
        "Camada 1 — data_loader.py   (ETL + Qualidade de Dados)",
        "  Leitura multi-aba Excel · fallback temporal · tratamento de S.I.",
        "Camada 2 — map_builder.py   (Algoritmos Geoespaciais + Visualização)",
        "  Algoritmo RAAE · construção do mapa Folium · camadas de patologia",
        "Camada 3 — app.py            (Interface do Dashboard + KPIs)",
        "  Interface Streamlit · filtros hierárquicos · indicadores de condição",
        "Camada 4 — report_generator.py   (Dado-para-Documento)",
        "  Relatório HTML automatizado · recomendações normativas de intervenção",
        "Stack: Python · Pandas · Folium/Leaflet.js · Streamlit · GeoPy · openpyxl",
    ])

    # 17
    _slide_conteudo(prs, "Algoritmo RAAE — Definição do Problema", [
        "Problema: partição espacial 1D ao longo do corredor ferroviário",
        "Atribuir raio r_i a cada OAE tal que:",
        "  Tangência mútua: r_i + r_j = d_ij  para OAEs vizinhas na AGM",
        "  Sem sobreposição entre OAEs não-adjacentes",
        "  Sem raios nulos ou negativos",
        "Equivalente a: decomposição de Voronoi 1D ao longo do eixo ferroviário",
        "  (Aurenhammer, 1991) — sem requerer a geometria do traçado como entrada",
        "Solução: algoritmo em 3 etapas — Union-Find + AGM (Prim) + Iteração de Jacobi",
    ])

    # 18
    _slide_duas_colunas(prs, "Algoritmo RAAE — Etapas 1 e 2",
        ["Etapa 1: Agrupamento por ε-Proximidade\n(Union-Find)",
         "Agrupa OAEs dentro de ε = 500 m de distância geodésica",
         "Union-Find com compressão de caminho",
         "Complexidade: O(n·α(n)) ≈ O(n) amortizado",
         "25 OAEs → 19 grupos representativos",
         "Trata estruturas co-localizadas (pontes gêmeas)",
         "Centroide do grupo = posição representativa",
         "Garante tangência visual na escala do grupo"],
        ["Etapa 2: Árvore Geradora Mínima\n(Algoritmo de Prim)",
         "AGM sobre k grupos, fila de prioridade",
         "Complexidade: O(k² log k)",
         "19 grupos → AGM com 18 arestas",
         "Captura topologia linear da ferrovia",
         "Sem geometria do traçado como entrada",
         "Adjacência AGM define restrições de tangência",
         "Equivalente funcional de Voronoi 1D"])

    # 19
    _slide_conteudo(prs, "Algoritmo RAAE — Etapa 3: Iteração de Jacobi", [
        "Sistema de restrições de tangência:",
        "  { r_i + r_j = d_ij  :  (i, j) ∈ AGM }  — uma equação por aresta",
        "Regra de atualização por relaxação simultânea (Equação 1):",
        "  r_i^(t+1)  =  max( r_min,  min_{j ∈ N(i)} ( d_ij − r_j^(t) ) )",
        "  onde  r_min = 800 m  |  buffer de 3% para tangência visualmente perceptível",
        "Convergência: max|Δr| < 1 m  (< 86 iterações para todos os grupos)",
        "Detecção de gap: d_ij − r_i − r_j > 1 km → linha tracejada cinza no mapa",
        "  Identifica trechos ferroviários sem cobertura de monitoramento",
        "Distâncias calculadas geodesicamente via GeoPy (WGS84) — precisão métrica",
    ])

    # ─── SEÇÃO 4: RESULTADOS ─────────────────────────────────────────────────
    _slide_divisor(prs, 4, "Resultados")

    # 21
    _slide_conteudo(prs, "Resultados ETL — Distribuição de Condição", [
        "25 OAEs consolidadas ao longo de 5 anos de inspeção (2021–2025)",
        "Fallback temporal: 14 OAEs de 2025 · 7 de 2024 · 4 de 2023",
        "Nenhuma OAE precisou de fallback além de 2023 (cobertura consistente)",
        "Distribuição de condição geral (G):",
        "  0 Emergencial · 0 Crítico · 4 Ruim · 13 Regular · 1 Bom · 0 Excelente",
        "  7 OAEs: S.I. — ciclo de inspeção especial pendente",
        "Destaque: 68% das OAEs em Regular (3) — horizonte de intervenção: 180 dias",
        "4 OAEs em Ruim (2) — intervenção necessária em até 90 dias",
    ])

    # 22
    _slide_tabela(prs, "Tabela 1 — Distribuição de Zonas RAAE",
        ["Nível RAAE", "Qtd.", "% do Total", "Cor no Mapa"],
        [
            ["Baixo",      "7",  "28%", "Verde"],
            ["Médio",      "12", "48%", "Amarelo"],
            ["Alto",       "4",  "16%", "Laranja"],
            ["Muito Alto", "2",  "8%",  "Vermelho"],
            ["Extremo",    "0",  "0%",  "Roxo"],
        ],
        nota="Tabela 1. Distribuição RAAE — 25 OAEs, corredor MG-ES. Fonte: O autor (2026).")

    # 23 ── Figura 2
    _slide_figura(prs,
        "Figura 2 — Zonas RAAE: Corredor MG-ES (905 km)",
        "Zonas mutuamente tangentes geradas por Union-Find + AGM (Prim) + Jacobi. "
        "Linhas tracejadas: 2 trechos sem cobertura (35,9 km e 171,2 km). Base: CartoDB Positron.",
        "CAPTURA: tile CartoDB Positron | Coloração = Por RAAE | camada Zonas RAAE ativa\n"
        "Zoom: todo o corredor visível | PNG mínimo 1400×900 px"
    )

    # 24 ── Figura 3
    _slide_figura(prs,
        "Figura 3 — Dashboard Interativo de Gestão à Vista",
        "KPIs de condição (NBR 9452:2019), filtros hierárquicos, mapa multicamada "
        "com zonas RAAE e legenda com contagem por nível.",
        "CAPTURA: tela completa do navegador | tile CartoDB Positron | Coloração = Por RAAE\n"
        "Ativar 1 camada de patologia | Legenda RAAE visível | PNG mínimo 1280×900 px"
    )

    # 25
    _slide_conteudo(prs, "Camadas de Manifestações Patológicas", [
        "10 categorias mapeadas a cores distintas (paleta de 17 cores):",
        "  Concreto: armadura exposta · fissuração · carbonatação · segregação · deformação · calcinação",
        "  Aço: corrosão · fratura · deformação · falha em ligação",
        "Cada categoria → camada Folium independente (LayerControl controlado pelo usuário)",
        "Leitura espacial habilitada pelo dashboard:",
        "  Quais tipos se concentram em quais km do corredor?",
        "  Há correlação entre zona RAAE e tipo de patologia predominante?",
        "Inteligência espacial qualitativa impossível de extrair apenas de dados tabulares",
    ])

    # 26
    _slide_conteudo(prs, "Relatório Técnico Automatizado — Dado-para-Documento", [
        "Usuário seleciona parâmetros de condição (ex.: 2 = Ruim + 3 = Regular)",
        "Sistema gera relatório HTML em < 2 segundos para 17 OAEs:",
        "  Capa: data de emissão, critérios selecionados, referência normativa",
        "  Resumo Executivo: distribuição de condição + distribuição RAAE",
        "  Fichas Individuais por OAE:",
        "    Dados estruturais (tipo, km, ano de construção, vão)",
        "    Resultados das inspeções rotineiras e especiais",
        "    Manifestações patológicas (texto + categoria)",
        "    Recomendações de intervenção vinculadas à NBR 9452 / NBR 6118 / NBR 12655",
        "Exportação para PDF via Ctrl+P do navegador (CSS de impressão otimizado)",
    ])

    # 27
    _slide_conteudo(prs, "Desempenho do Algoritmo e do Sistema", [
        "Pipeline ETL: < 3 segundos para 25 OAEs × 5 anos (125 registros)",
        "Agrupamento Union-Find: 25 → 19 grupos · 6 pares co-localizados identificados",
        "AGM de Prim: 18 arestas sobre 19 nós",
        "Convergência do Jacobi: máximo 86 iterações para todos os grupos",
        "Tangência RAAE: 16 / 18 arestas (88,9%) plenamente tangentes",
        "Detecção de gap: 2 trechos não monitorados identificados (35,9 km + 171,2 km)",
        "Geração de relatório: 17 OAEs em relatório HTML em < 2 segundos",
        "Escalabilidade: O(n² log n) — 389 OAEs processáveis em segundos em hardware convencional",
    ])

    # ─── SEÇÃO 5: CONCLUSÕES ─────────────────────────────────────────────────
    _slide_divisor(prs, 5, "Conclusões")

    # 29
    _slide_conteudo(prs, "Principais Contribuições", [
        "(1) Pipeline ETL multitemporal",
        "  Lógica de fallback + sentinela S.I. para registros de inspeção ausentes",
        "(2) Algoritmo original de delimitação de zonas RAAE",
        "  Union-Find + AGM (Prim) + Iteração de Jacobi → Voronoi 1D equivalente",
        "  Sem geometria do traçado ferroviário como entrada",
        "(3) Dashboard interativo multicamada",
        "  25 OAEs · 905 km · KPIs de condição · camadas de patologia · legenda RAAE",
        "(4) Gerador de relatório técnico automatizado",
        "  Recomendações NBR 9452:2019 · geração em < 2 segundos",
    ])

    # 30
    _slide_duas_colunas(prs, "Limitações e Trabalhos Futuros",
        ["Limitações",
         "Piloto: 25 de 389 OAEs — poder estatístico limitado",
         "PLN léxico: sem desambiguação contextual para patologias compostas",
         "Jacobi: sem garantia de otimalidade global (estável para grafos quase-lineares)",
         "Sem integração de registros fotográficos ainda",
         "Relatório: HTML/PDF apenas — sem integração com ERP de manutenção"],
        ["Trabalhos Futuros",
         "Fase 2: expandir para 389 OAEs, corredor completo de 905 km",
         "Legendagem de imagens de inspeção (Brilakis et al., 2023)",
         "PLN semântico (Zheng et al., 2022) para Fase 2",
         "LSTM / Transformer para previsão de degradação estrutural",
         "Dado de treino: banco de inspeção multianual estruturado por este pipeline"])

    # 31
    _slide_conteudo(prs, "Conclusões", [
        "Ferramentas de código aberto + teoria dos grafos RESOLVEM a fragmentação",
        "  informacional dos sistemas de inspeção de OAEs ferroviárias",
        "A abordagem data-centric mostrou-se mais impactante que complexidade algorítmica:",
        "  Estruturar 5 anos de dados heterogêneos = contribuição principal do trabalho",
        "A detecção de gaps RAAE transforma lacuna de dados em recomendação de gestão:",
        "  → 35,9 km e 171,2 km identificados como trechos sem monitoramento estrutural",
        "Base estabelecida para Fase 2 (389 OAEs) e modelos preditivos DNNs",
        "Todas as ferramentas: código aberto (Streamlit · Folium · Pandas · GeoPy · Python)",
        "  Custo zero de licenciamento — implantável por qualquer operadora ferroviária",
    ])

    # 32
    _slide_conteudo(prs, "Referências (APA 7ª edição)", [
        "Aurenhammer, F. (1991). Voronoi diagrams. ACM Computing Surveys, 23(3), 345–405.",
        "Brilakis et al. (2023). Image captioning for bridge inspection. Automation in Constr., 148.",
        "Kang et al. (2020). Deep NN para extração de informação. Automation in Constr., 110.",
        "Li et al. (2023). Few-shot learning para inspeção de pontes. J. Bridge Eng., 28(4).",
        "McKinney, W. (2010). Data structures for statistical computing in Python. SciPy Proc.",
        "Ng, A. (2021). From model-centric to data-centric AI [Vídeo]. DeepLearning.AI.",
        "Pozzi et al. (2017). Bridge management with BIM/GIS. J. Bridge Eng., 22(11).",
        "Zheng et al. (2022). PLN semântico para inspeção de pontes. CACAIE, 37(8).",
        "ABNT NBR 9452:2019 — Inspeção de pontes, viadutos e passagens subterrâneas.",
        "ABNT NBR 6118:2014 — Projeto de estruturas de concreto.",
    ])

    # 33 ── Encerramento
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_bg(sl, '1F3864')
    _rect(sl, 0, 3.05, 10, 0.07, '2E75B6')
    _box(sl, 1, 1.1, 8, 1.5,
         "Obrigado!", fs=52, bold=True, color=P_WHITE, align=PP_ALIGN.CENTER)
    _box(sl, 1, 3.25, 8, 0.75,
         "Perguntas e Discussão", fs=26, italic=True,
         color=P_LBLUE, align=PP_ALIGN.CENTER)
    _box(sl, 1, 4.20, 8, 1.05,
         "Wagner de Almeida Tavares\n"
         "Instituto Militar de Engenharia (IME)\n"
         "wagner.tavares2024@gmail.com",
         fs=14, color=P_LBLUE, align=PP_ALIGN.CENTER)
    _rodape(sl)

    out = os.path.join(OUT_DIR, "apresentacao_oae_ET261400_ptBR.pptx")
    prs.save(out)
    print(f"[PPTX] Salvo: {out}")
    print(f"[PPTX] Total de slides: {len(prs.slides)}")
    return out


# ════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("\nConcluído. Abra os arquivos na pasta do projeto.")
